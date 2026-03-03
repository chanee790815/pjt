import streamlit as st
import os
import io
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime

COLOR_MAIN_PPT = RGBColor(16, 185, 129)

def get_clean_text(element):
    if not element: return ""
    return element.get_text(separator=' ', strip=True).replace('  ', ' ')

def fill_placeholders(slide, title_text):
    if slide.shapes.title: slide.shapes.title.text = title_text
    placeholders = [p for p in slide.placeholders if p.placeholder_format.idx != 0]
    if placeholders:
        sub_p = min(placeholders, key=lambda p: p.top)
        sub_p.text = "2. 개발사업부_PM팀"
        for p in sub_p.text_frame.paragraphs:
            p.font.name = '맑은 고딕'; p.font.size = Pt(20); p.font.bold = True

def create_card(slide, left, top, width, height, title, body, is_kpi=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 255, 255) if is_kpi else RGBColor(248, 250, 252)
    shape.line.color.rgb = COLOR_MAIN_PPT if is_kpi else RGBColor(220, 227, 235)
    t_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.5))
    p1 = t_box.text_frame.paragraphs[0]; p1.text = title; p1.font.bold = True; p1.font.name = '맑은 고딕'; p1.font.size = Pt(11 if is_kpi else 13)
    if is_kpi: p1.alignment = PP_ALIGN.CENTER
    b_box = slide.shapes.add_textbox(left + Inches(0.1), top + (height * 0.45 if is_kpi else Inches(0.7)), width - Inches(0.2), height * 0.5)
    p2 = b_box.text_frame.paragraphs[0]; p2.text = body.replace('\n', ' ').strip(); p2.font.name = '맑은 고딕'
    if is_kpi:
        p2.font.bold = True; p2.font.color.rgb = COLOR_MAIN_PPT; p2.alignment = PP_ALIGN.CENTER; p2.font.size = Pt(18)
    else:
        p2.font.size = Pt(11)

def view_ppt_generator(DEFAULT_TEMPLATE):
    st.title("📊 워크샵 PPT 자동 생성기")
    t_exist = os.path.exists(DEFAULT_TEMPLATE)
    if t_exist: st.success(f"✅ 서버 양식 로드 완료: `{DEFAULT_TEMPLATE}`")
    else: st.warning("⚠️ 서버에 양식 파일이 없습니다.")
    h_up = st.file_uploader("데이터 파일 업로드 (input.html)", type=['html'])
    p_up = st.file_uploader("양식 교체 (선택사항)", type=['pptx'])
    
    if st.button("🚀 PPT 즉시 생성 시작"):
        f_temp = p_up if p_up else (DEFAULT_TEMPLATE if t_exist else None)
        if h_up and f_temp:
            try:
                soup = BeautifulSoup(h_up, 'lxml'); prs = Presentation(f_temp if not p_up else io.BytesIO(p_up.read()))
                for s in list(prs.slides._sldIdLst): prs.slides._sldIdLst.remove(s)
                W, H, L5 = prs.slide_width, prs.slide_height, prs.slide_layouts[4]
                
                # Slide 1: 표지
                s1 = soup.select_one("#slide1")
                if s1:
                    sld = prs.slides.add_slide(prs.slide_layouts[0])
                    for sh in sld.placeholders:
                        if sh.placeholder_format.type == 1: sh.text = get_clean_text(s1.find('h1'))
                        elif sh.placeholder_format.type == 2: sh.text = get_clean_text(s1.find('p'))
                
                # Slide 2: KPI (8개 그리드)
                s2 = soup.select_one("#slide2")
                if s2:
                    sld = prs.slides.add_slide(L5); fill_placeholders(sld, "1. 2026년 PM팀 핵심 성과지표 (KPI)")
                    cards = s2.select(".kpi-card"); card_h = H * 0.18
                    for i, card in enumerate(cards):
                        r, c = i // 4, i % 4
                        create_card(sld, W*0.06+(c*(W*0.225)), H*0.3+(r*(card_h+Inches(0.1))), W*0.21, card_h, get_clean_text(card.find(class_="label")), get_clean_text(card.find(class_="val")), is_kpi=True)
                
                # Slide 3: 핵심 역량
                s3 = soup.select_one("#slide3")
                if s3:
                    sld = prs.slides.add_slide(L5); fill_placeholders(sld, "2. 핵심 역량: 계통 인프라 기술 자산화")
                    tb = sld.shapes.add_textbox(W*0.06, H*0.3, W*0.55, H*0.6)
                    for li in s3.select("li"): p = tb.text_frame.add_paragraph(); p.text = "• " + get_clean_text(li); p.font.size = Pt(17); p.space_after = Pt(12)
                
                # Slide 4: AI 전략
                s4 = soup.select_one("#slide4")
                if s4:
                    sld = prs.slides.add_slide(L5); fill_placeholders(sld, "3. AI 기반 에너지 락인(Lock-in) 전략")
                    rect = sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, W*0.06, H*0.3, W*0.65, H*0.55); rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(240, 253, 244); rect.line.color.rgb = COLOR_MAIN_PPT
                    tb = sld.shapes.add_textbox(W*0.08, H*0.33, W*0.6, H*0.5)
                    for item in s4.select(".ai-tag, h3, li, p"): 
                        text = get_clean_text(item)
                        if text: p = tb.text_frame.add_paragraph(); p.text = text; p.font.size = Pt(16); p.space_after = Pt(10)
                
                # Slide 5: R&R
                s5 = soup.select_one("#slide5")
                if s5:
                    sld = prs.slides.add_slide(L5); fill_placeholders(sld, "4. PM팀 핵심 실무 그룹 R&R")
                    tiles = s5.select(".tile"); cols = 3; card_w = (W*0.88)/cols; card_h = H*0.25
                    for i, tile in enumerate(tiles): r, c = i // cols, i % cols; create_card(sld, W*0.06+(c*(card_w+Inches(0.1))), H*0.3+(r*(card_h+Inches(0.1))), card_w, card_h, get_clean_text(tile.find("h3")), get_clean_text(tile.find("p")))
                
                # Slide 6: 인원 충원 (인원수 동적 추출)
                s6 = soup.select_one("#slide6")
                if s6:
                    sld = prs.slides.add_slide(L5); fill_placeholders(sld, "5. 미래 파이프라인 인원 충원 계획")
                    num_val = get_clean_text(s6.find("div", style=lambda x: x and "font-size: 100px" in x))
                    num_tb = sld.shapes.add_textbox(W*0.06, H*0.45, W*0.35, Inches(1.5)); p = num_tb.text_frame.paragraphs[0]; p.text = f"{num_val} 명"; p.font.size = Pt(84); p.font.bold = True; p.font.color.rgb = COLOR_MAIN_PPT; p.alignment = PP_ALIGN.CENTER
                    dtb = sld.shapes.add_textbox(W*0.42, H*0.3, W*0.52, H*0.6)
                    for p_t in s6.select("p"): text = get_clean_text(p_t); 
                    if text: p = dtb.text_frame.add_paragraph(); p.text = text; p.font.size = Pt(17); p.space_after = Pt(12)
                
                # Slide 7: Vision
                s7 = soup.select_one("#slide7")
                if s7:
                    sld = prs.slides.add_slide(L5); fill_placeholders(sld, "Vision 2026")
                    tb = sld.shapes.add_textbox(W*0.1, H*0.5, W*0.8, Inches(1.5)); p = tb.text_frame.paragraphs[0]; p.text = get_clean_text(s7.find("p")); p.font.size = Pt(28); p.font.bold = True; p.alignment = PP_ALIGN.CENTER

                ppt_out = io.BytesIO(); prs.save(ppt_out)
                st.success("✅ 생성 완료!"); st.download_button("📥 다운로드", ppt_out.getvalue(), file_name=f"신성_PM전략_{datetime.date.today()}.pptx")
            except Exception as e: st.error(f"오류: {e}")
